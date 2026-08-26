# mypy: disable-error-code="no-untyped-call,attr-defined,import-untyped,valid-type"
#
# Neither `python-pptx` nor `reportlab` ships type stubs, and pptx's
# `Presentation()` is an untyped factory rather than a class, so it cannot
# even be used as an annotation. The relaxation is scoped to this one
# module rather than loosened globally — every other file in the package
# stays under full --strict.
"""PDF and PowerPoint export — M10 section 5.6.

The exit criterion for the whole project is that a complete scenario
produces a distributable, fully cited deliverable. This module is that
final step: it takes a calculated result plus its narrative and emits a
file someone can send to a colleague.

Both formats carry the full assumption register and the mandatory
limitations. That is not padding — a budget impact number without its
assumptions is not reviewable, and the register is what makes the estimate
defensible rather than merely presentable.
"""

from __future__ import annotations

import io
from datetime import UTC, datetime

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.platypus import (
    KeepTogether,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

from ..schemas.calculation import CalculationResponse, EvidenceGapResponse
from .narrative_service import Narrative

#: Matches the interface and the published run report, so a exported deck
#: and the screen it came from read as one product.
ACCENT = colors.HexColor("#07707C")
INK = colors.HexColor("#0E181B")
INK_2 = colors.HexColor("#42585F")
LINE = colors.HexColor("#D8E3E6")
SURFACE_2 = colors.HexColor("#EDF3F4")

SECTION_TITLES = {
    "population": "Population",
    "impact": "Budget impact",
    "affordability": "Affordability",
    "uncertainty": "Uncertainty",
    "limitations": "Limitations",
}


def _money(amount: float, currency: str) -> str:
    if abs(amount) >= 1e9:
        return f"{amount / 1e9:,.2f} bn {currency}"
    if abs(amount) >= 1e6:
        return f"{amount / 1e6:,.1f} m {currency}"
    return f"{amount:,.0f} {currency}"


# --------------------------------------------------------------------------- PDF


def build_pdf(
    result: CalculationResponse,
    narrative: Narrative,
    asset: str,
    gaps: EvidenceGapResponse | None = None,
) -> bytes:
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer, pagesize=A4,
        leftMargin=20 * mm, rightMargin=20 * mm,
        topMargin=18 * mm, bottomMargin=18 * mm,
        title=f"Budget impact — {asset}",
        author="BIET",
    )

    base = getSampleStyleSheet()
    h1 = ParagraphStyle("h1", parent=base["Title"], fontSize=19, leading=23,
                        textColor=INK, alignment=TA_LEFT, spaceAfter=2)
    sub = ParagraphStyle("sub", parent=base["Normal"], fontSize=9.5, leading=13,
                         textColor=INK_2, spaceAfter=14)
    h2 = ParagraphStyle("h2", parent=base["Heading2"], fontSize=10, leading=13,
                        textColor=ACCENT, spaceBefore=15, spaceAfter=6)
    body = ParagraphStyle("body", parent=base["Normal"], fontSize=9.5, leading=14,
                          textColor=INK, spaceAfter=7)
    small = ParagraphStyle("small", parent=base["Normal"], fontSize=8, leading=11,
                           textColor=INK_2)
    figure = ParagraphStyle("figure", parent=base["Normal"], fontSize=25, leading=29,
                            textColor=INK, spaceAfter=3)

    story: list[object] = []

    story.append(Paragraph(f"Budget impact — {asset}", h1))
    story.append(Paragraph(
        f"{len(result.countries)} markets · launch {result.launch_year} · "
        f"{result.horizon_years}-year horizon · reported in {result.reporting_currency} · "
        f"engine {result.engine_version} · FX {result.fx_snapshot_date} · "
        f"generated {datetime.now(UTC):%Y-%m-%d} by {narrative.generated_by}",
        sub,
    ))

    story.append(Paragraph("Cumulative incremental budget impact", h2))
    story.append(Paragraph(
        _money(result.totals.cumulative, result.totals.currency), figure,
    ))
    story.append(Paragraph(
        "Incremental — the world with this asset minus the world without it, net of "
        f"displaced therapy. Peak in year {result.totals.peak_year}.", small,
    ))
    story.append(Spacer(1, 12))

    # --- per-market table
    rows = [["Market", "Addressable", "On therapy", "Cumulative impact",
             "Price basis", "Affordability"]]
    for country in result.countries:
        last = country.years[-1]
        rows.append([
            f"{country.country_code} ({country.currency})",
            f"{last.addressable:,.0f}",
            f"{last.patients_on_new:,.0f}",
            f"{country.cumulative_budget_impact:,.0f}",
            country.new_therapy.price_basis.replace("_", " "),
            (f"{country.affordability.band} · "
             f"{country.affordability.cumulative_ratio * 100:.3f}%"
             if country.affordability else "—"),
        ])
    table = Table(rows, hAlign="LEFT", colWidths=[68, 68, 62, 88, 62, 82])
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), SURFACE_2),
        ("TEXTCOLOR", (0, 0), (-1, 0), INK_2),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 7.5),
        ("ALIGN", (1, 0), (-1, -1), "RIGHT"),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("GRID", (0, 0), (-1, -1), 0.4, LINE),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
    ]))
    story.append(table)

    # --- narrative
    for key, title in SECTION_TITLES.items():
        if key in narrative.sections:
            story.append(KeepTogether([
                Paragraph(title, h2),
                Paragraph(narrative.sections[key], body),
            ]))

    # --- citations
    if narrative.citations:
        story.append(Paragraph("Cited guidance", h2))
        for chunk in narrative.citations:
            story.append(Paragraph(
                f"{chunk.issuing_body} — {chunk.document_title}, "
                f"p.{chunk.page_number} (similarity {chunk.similarity:.2f})", small,
            ))

    # --- limitations
    story.append(Paragraph("Stated limitations", h2))
    for limitation in narrative.limitations:
        story.append(Paragraph(f"• {limitation}", small))
        story.append(Spacer(1, 3))

    # --- evidence priorities (M15), before the register rather than after:
    # the register says what every value rests on, and this says which of
    # those are worth doing something about. A reader who stops early should
    # have seen the shorter, more actionable list.
    if gaps is not None and gaps.gaps:
        story.append(Spacer(1, 10))
        story.append(Paragraph("What to find out next", h2))
        story.append(Paragraph(
            "Ranked by how much each assumption moves this result multiplied by how "
            "weakly it is founded. A large swing on a published country-specific source "
            "is settled; a large swing on a placeholder is why this answer cannot yet be "
            "relied on. A parameter that cannot move the result is never a priority, "
            "however weak its source.",
            sub,
        ))
        priorities = [["Priority", "Parameter", "Moves", "Tier", "Rests on"]]
        for gap in gaps.gaps:
            priorities.append([
                gap.priority.upper(),
                gap.label,
                _money(gap.swing, gaps.currency),
                gap.confidence_tier,
                (gap.source[:70] + "…") if len(gap.source) > 70 else gap.source,
            ])
        gap_table = Table(
            priorities, hAlign="LEFT", colWidths=[52, 105, 62, 26, 185], repeatRows=1,
        )
        gap_table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), SURFACE_2),
            ("TEXTCOLOR", (0, 0), (-1, 0), INK_2),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 6.6),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("GRID", (0, 0), (-1, -1), 0.3, LINE),
            ("TOPPADDING", (0, 0), (-1, -1), 3),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
        ]))
        story.append(gap_table)

    # --- assumption register, on its own page: it is the audit trail, and
    # burying it under a table would be the wrong signal about its status.
    story.append(PageBreak())
    story.append(Paragraph("Assumption register", h1))
    story.append(Paragraph(
        "Every resolved value this run consumed, with the source it came from and how "
        "much weight it carries. Built from the run's own stored snapshot, so this "
        "reflects what the run actually used rather than what the database holds today.",
        sub,
    ))

    register = [["Parameter", "Market", "Value", "Tier", "Source"]]
    for entry in narrative.assumptions:
        register.append([
            entry.parameter_path,
            entry.country_code or "all",
            f"{entry.value:,.4g}",
            str(entry.confidence_tier),
            (entry.source[:78] + "…") if len(entry.source) > 78 else entry.source,
        ])
    reg_table = Table(register, hAlign="LEFT", colWidths=[95, 34, 52, 26, 223], repeatRows=1)
    reg_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), SURFACE_2),
        ("TEXTCOLOR", (0, 0), (-1, 0), INK_2),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 6.4),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("GRID", (0, 0), (-1, -1), 0.3, LINE),
        ("TOPPADDING", (0, 0), (-1, -1), 3),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
    ]))
    story.append(reg_table)

    doc.build(story)
    return buffer.getvalue()


# --------------------------------------------------------------------------- PPTX


def _text_slide(prs: object, title: str, body: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    frame = slide.placeholders[1].text_frame
    frame.text = body
    frame.word_wrap = True
    for para in frame.paragraphs:
        para.font.size = Pt(15)


def build_pptx(
    result: CalculationResponse,
    narrative: Narrative,
    asset: str,
    gaps: EvidenceGapResponse | None = None,
) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)          # 16:9, the shape a deck is shown in
    prs.slide_height = Inches(7.5)

    # --- title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"Budget impact — {asset}"
    slide.placeholders[1].text = (
        f"{len(result.countries)} markets · launch {result.launch_year} · "
        f"{result.horizon_years}-year horizon\n"
        f"Reported in {result.reporting_currency} · engine {result.engine_version} · "
        f"FX {result.fx_snapshot_date}"
    )

    # --- headline
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Cumulative incremental budget impact"
    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.6))
    frame = box.text_frame
    frame.word_wrap = True
    run = frame.paragraphs[0].add_run()
    run.text = _money(result.totals.cumulative, result.totals.currency)
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x07, 0x70, 0x7C)

    note = frame.add_paragraph()
    note.text = (
        "Incremental — the world with this asset minus the world without it, "
        f"net of displaced therapy. Peak in year {result.totals.peak_year}."
    )
    note.font.size = Pt(14)
    note.font.color.rgb = RGBColor(0x42, 0x58, 0x5F)

    # --- per-market table
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "By market"
    rows, cols = len(result.countries) + 1, 5
    shape = slide.shapes.add_table(
        rows, cols, Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.4 * rows),
    )
    table = shape.table
    for i, heading in enumerate(
        ["Market", "Addressable", "Cumulative impact", "Price basis", "Affordability"]
    ):
        table.cell(0, i).text = heading
    for r, country in enumerate(result.countries, start=1):
        last = country.years[-1]
        cells = [
            f"{country.country_code} ({country.currency})",
            f"{last.addressable:,.0f}",
            f"{country.cumulative_budget_impact:,.0f}",
            country.new_therapy.price_basis.replace("_", " "),
            (f"{country.affordability.band} · "
             f"{country.affordability.cumulative_ratio * 100:.3f}%"
             if country.affordability else "—"),
        ]
        for i, value in enumerate(cells):
            table.cell(r, i).text = value
    for row in table.rows:
        for cell in row.cells:
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(11)

    # --- narrative, one slide per section
    for key, title in SECTION_TITLES.items():
        if key in narrative.sections:
            _text_slide(prs, title, narrative.sections[key])

    # --- limitations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Stated limitations"
    frame = slide.placeholders[1].text_frame
    frame.word_wrap = True
    frame.text = narrative.limitations[0]
    for limitation in narrative.limitations[1:]:
        para = frame.add_paragraph()
        para.text = limitation
    for para in frame.paragraphs:
        para.font.size = Pt(11)

    # --- citations
    if narrative.citations:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Cited guidance"
        frame = slide.placeholders[1].text_frame
        frame.word_wrap = True
        frame.text = ""
        for chunk in narrative.citations:
            para = frame.add_paragraph()
            para.text = (
                f"{chunk.issuing_body} — {chunk.document_title}, p.{chunk.page_number}"
            )
            para.font.size = Pt(12)

    # --- evidence priorities (M15). One slide, actionable list only: a deck
    # is where a decision gets argued, and "settled" rows are not the argument.
    if gaps is not None:
        actionable = [g for g in gaps.gaps if g.priority != "sufficient"]
        if actionable:
            _text_slide(
                prs, "What to find out next",
                "\n".join(
                    f"{g.priority.upper()} — {g.label}: moves "
                    f"{_money(g.swing, gaps.currency)}, tier {g.confidence_tier}. "
                    f"Rests on {g.source[:90]}"
                    for g in actionable
                ),
            )

    # --- assumption register
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Assumption register"
    frame = slide.placeholders[1].text_frame
    frame.word_wrap = True
    frame.text = (
        f"{len(narrative.assumptions)} resolved values, each with its source and "
        "confidence tier. Full detail in the PDF export."
    )
    for entry in narrative.assumptions[:14]:
        para = frame.add_paragraph()
        para.text = (
            f"{entry.parameter_path} ({entry.country_code or 'all'}): "
            f"{entry.value:,.4g} — tier {entry.confidence_tier}"
        )
        para.font.size = Pt(10)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
